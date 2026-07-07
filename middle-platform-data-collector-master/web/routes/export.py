"""数据导出 API"""
import base64
import io
import logging
import os
import re
import shutil
import tempfile
import time
import uuid
from datetime import datetime
from pathlib import Path
from io import BytesIO

from flask import Blueprint, request, jsonify, send_file, session

from models.weekly_record import WeeklyRecord
from models.monthly_record import MonthlyRecord
from services.exporter import export_weekly, export_monthly

logger = logging.getLogger(__name__)
export_bp = Blueprint("export", __name__)

# PPT 会话存储：{session_id: {"dir": Path, "start_date": str, "end_date": str, "school_name": str, "created_at": float}}
_PPT_SESSIONS = {}
_PPT_SESSION_TTL = 600  # 10 分钟过期


def _safe_filename(label: str) -> str:
    return re.sub(r'[\\/:*?"<>|\s]', '_', label)


def _get_allowed_schools():
    """获取当前用户可见的学校列表"""
    user_id = session.get("user_id")
    if not user_id:
        return None
    from models.user import User
    user = User.get_by_id(user_id)
    if not user:
        return None
    if user.is_admin or not user.assigned_schools:
        return None
    return user.school_list


@export_bp.route("/weekly")
def export_weekly_data():
    """导出单周数据"""
    allowed = _get_allowed_schools()
    year = request.args.get("year", type=int)
    week_number = request.args.get("week_number", type=str) or ""
    school_name = request.args.get("school_name", "")
    month_prefix = request.args.get("month_prefix", "")

    if not year:
        return jsonify({"error": "请提供 year 参数"}), 400

    # 非管理员如果选择了不在自己范围内的学校，拒绝
    if allowed is not None and school_name and school_name not in allowed:
        return jsonify({"error": "无权访问该学校数据"}), 403

    records = WeeklyRecord.query_flexible(year, week_number, school_name, month_prefix)
    # 过滤记录
    if allowed is not None:
        records = [r for r in records if r.school_name in allowed]
    if not records:
        return jsonify({"error": "没有找到对应数据"}), 404

    filepath = export_weekly(records, year, week_number)
    from services.exporter import _build_export_name
    _period = week_number or (month_prefix if month_prefix else "")
    dl_name = _build_export_name(year, _period) + ".xlsx"
    return send_file(
        filepath,
        as_attachment=True,
        download_name=dl_name,
    )


@export_bp.route("/preview")
def preview_data():
    """预览数据（JSON）"""
    allowed = _get_allowed_schools()
    year = request.args.get("year", type=int)
    week_number = request.args.get("week_number", type=str) or ""
    school_name = request.args.get("school_name", "")
    month_prefix = request.args.get("month_prefix", "")

    if not year:
        return jsonify({"error": "请提供 year 参数"}), 400

    # 非管理员如果选择了不在自己范围内的学校，返回空
    if allowed is not None and school_name and school_name not in allowed:
        return jsonify({"records": []})

    records = WeeklyRecord.query_flexible(year, week_number, school_name, month_prefix)
    if allowed is not None:
        records = [r for r in records if r.school_name in allowed]
    return jsonify({"records": [r.to_dict() for r in records]})


@export_bp.route("/monthly")
def export_monthly_data():
    """导出月度数据"""
    allowed = _get_allowed_schools()
    year = request.args.get("year", type=int)
    month_number = request.args.get("month_number", "")
    school_name = request.args.get("school_name", "")

    if not year:
        return jsonify({"error": "请提供 year 参数"}), 400

    if allowed is not None and school_name and school_name not in allowed:
        return jsonify({"error": "无权访问该学校数据"}), 403

    records = MonthlyRecord.query_flexible(year, month_number, school_name)
    if allowed is not None:
        records = [r for r in records if r.school_name in allowed]
    if not records:
        return jsonify({"error": "没有找到对应数据"}), 404

    filepath = export_monthly(records, year, month_number)
    from services.exporter import _build_export_name
    dl_name = _build_export_name(year, month_number) + ".xlsx"
    return send_file(
        filepath,
        as_attachment=True,
        download_name=dl_name,
    )

@export_bp.route("/distinct_weeks")
def distinct_weeks():
    """查询指定年份已有的不重复周标签"""
    year = request.args.get("year", type=int)
    if not year:
        return jsonify({"error": "请提供 year 参数"}), 400
    weeks = WeeklyRecord.query_distinct_weeks(year)
    return jsonify({"weeks": weeks})


# ── PPT 模板常量 ──
_PPT_BLUE = "1A73E8"       # 天立主色
_PPT_DARK = "0D47A1"       # 深蓝
_PPT_LIGHT_BG = "F0F4FA"   # 浅蓝背景
_PPT_WHITE = "FFFFFF"

# PPT 内容说明文字
_PPT_DESC = {
    "kpi": "本页展示选定时间范围内全校总览的核心KPI指标，包括学校总数、日/周/月活跃人数及作业提交情况。",
    "bar": "各校使用率排名（柱状图），反映各直营校在平台上的整体活跃度差异。",
    "donut_radar": "左图：优先级分层占比（环形图），展示高/中/低优先级学校分布。右图：全板块均衡度分析（雷达图），覆盖六大维度。",
    "trend": "日活人数趋势（折线图），按学段拆分展示各学段每日活跃教师人数变化。",
    "school_kpi": "本页展示指定学校在选定时间范围内的核心KPI指标，包含平台使用率、日/周/月活跃人数及作业数据。",
    "school_trend": "本校周期使用率趋势，按学段（高中/初中/小学）拆分展示每日使用率变化。",
    "school_module": "本校8大业务板块使用率拆解（柱状图），展示各板块活跃教师占比。",
}


def _make_pptx(screenshots, start_date, end_date, school_name, ppt_title=""):
    """使用模板 PPT 生成报告文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    import os

    # 加载模板
    template_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "static", "template.pptx")
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"模板文件不存在: {template_path}")
    prs = Presentation(template_path)

    def _replace_shape_with_image(slide, shape_name, img_path):
        """替换指定名称的形状为图片"""
        if not img_path or not os.path.exists(img_path):
            return False
        for shape in slide.shapes:
            if shape.name == shape_name:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(img_path, left, top, width, height)
                return True
        return False

    def _update_shape_text(slide, shape_name, new_text):
        """更新指定形状的文字内容"""
        for shape in slide.shapes:
            if shape.name == shape_name:
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        para.text = new_text
                    return True
        return False

    slides = list(prs.slides)

    # ── 第1页: 封面 - 更新标题 ──
    if len(slides) >= 1:
        title_text = ppt_title if ppt_title else "教育平台数据看板"
        _update_shape_text(slides[0], "矩形 12", title_text)

    # ── 第2页: 章节页 - 更新时间范围 ──
    if len(slides) >= 2:
        _update_shape_text(slides[1], "矩形 7", f"{start_date} ~ {end_date}")

    # ── 第3页: 全校总览 KPI + 活跃比例对比 ──
    if len(slides) >= 3:
        # KPI 小卡片保留模板默认文字（展示KPI名称）
        # 大区域放置活跃比例对比图
        img_ar = screenshots.get("active_ratio")
        if img_ar and os.path.exists(img_ar):
            _replace_shape_with_image(slides[2], "圆角矩形 1", img_ar)
        elif screenshots.get("kpi") and os.path.exists(screenshots["kpi"]):
            _replace_shape_with_image(slides[2], "圆角矩形 1", screenshots["kpi"])

    # ── 第4页: 各校使用率排名 ──
    if len(slides) >= 4:
        img = screenshots.get("bar")
        if img:
            _replace_shape_with_image(slides[3], "圆角矩形 1", img)

    # ── 第5页: 各校作业次数统计 ──
    if len(slides) >= 5:
        img = screenshots.get("homework")
        if img:
            _replace_shape_with_image(slides[4], "圆角矩形 1", img)

    # ── 第6页: 章节分隔页 (本校) - 无需修改 ──

    # ── 第7页: 单校 KPI + 短板预警 ──
    if len(slides) >= 7 and school_name:
        # KPI 小卡片保留模板默认文字（展示KPI名称）
        # 大区域放置 KPI + 短板预警整块截图
        img_combined = screenshots.get("school_kpi_warning")
        if img_combined and os.path.exists(img_combined):
            _replace_shape_with_image(slides[6], "圆角矩形 1", img_combined)

    # ── 第8页: 趋势分析 tab1 (使用率趋势) ──
    if len(slides) >= 8 and school_name:
        img = screenshots.get("school_trend_usage")
        if img:
            _replace_shape_with_image(slides[7], "圆角矩形 1", img)

    # ── 第9页: 趋势分析 tab2 (活跃度趋势) ──
    if len(slides) >= 9 and school_name:
        img = screenshots.get("school_trend_activity")
        if img:
            _replace_shape_with_image(slides[8], "圆角矩形 1", img)

    # ── 第10页: 趋势分析 tab3 (作业趋势) ──
    if len(slides) >= 10 and school_name:
        img = screenshots.get("school_trend_homework")
        if img:
            _replace_shape_with_image(slides[9], "圆角矩形 1", img)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _capture_screenshots(base_url, start_date, end_date, school_id, school_name, dest_dir=None):
    """使用 Playwright 截取页面截图"""
    from playwright.sync_api import sync_playwright

    def _wait_for_canvas_rendered(page, canvas_id, timeout=10000):
        """等待canvas有实际内容（非空白）"""
        import time
        start = time.time()
        while time.time() - start < timeout / 1000:
            has_content = page.evaluate(f"""() => {{
                const canvas = document.getElementById('{canvas_id}');
                if (!canvas) return false;
                const ctx = canvas.getContext('2d');
                if (!ctx) return false;
                const data = ctx.getImageData(0, 0, canvas.width, canvas.height).data;
                // Check if canvas has any non-transparent, non-white pixels
                for (let i = 0; i < data.length; i += 4) {{
                    if (data[i+3] > 0 && (data[i] < 250 || data[i+1] < 250 || data[i+2] < 250)) {{
                        return true;
                    }}
                }}
                return false;
            }}""")
            if has_content:
                return True
            time.sleep(0.5)
        return False

    screenshots = {}
    if dest_dir:
        tmp_dir = Path(dest_dir)
    else:
        tmp_dir = Path(tempfile.mkdtemp(prefix="ppt_export_"))
    tmp_dir.mkdir(parents=True, exist_ok=True)

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(channel="chrome", headless=True, args=["--disable-gpu", "--no-sandbox", "--disable-extensions"])
            context = browser.new_context(viewport={"width": 1920, "height": 1080})
            page = context.new_page()

            # ── 0. 自动登录 ──
            logger.info("[PPT] 自动登录...")
            # 先访问登录页（无需认证），拿到同源上下文
            page.goto(f"{base_url}/login", wait_until="domcontentloaded", timeout=20000)
            login_result = page.evaluate("""
                async () => {
                    const resp = await fetch('/login', {
                        method: 'POST',
                        headers: {'Content-Type': 'application/x-www-form-urlencoded'},
                        body: 'username=13931822731&password=qiming123&next=%2F'
                    });
                    return await resp.json();
                }
            """)
            if not login_result.get("success"):
                logger.error("[PPT] 登录失败: %s", login_result)
                raise Exception(f"自动登录失败: {login_result.get('error', '未知错误')}")
            # 登录后跳转首页让 session 生效
            page.goto(base_url, wait_until="domcontentloaded", timeout=30000)
            time.sleep(2)  # 等待 JS 初始化
            logger.info("[PPT] 登录成功")

            # ── 1. 首页 KPI ──
            logger.info("[PPT] 采集首页: dates=%s ~ %s", start_date, end_date)
            page.goto(base_url, wait_until="load", timeout=30000)
            time.sleep(2)
            # 设置日期并触发查询
            page.evaluate(f"""() => {{
                document.getElementById('f_start').value = '{start_date}';
                document.getElementById('f_end').value = '{end_date}';
            }}""")
            time.sleep(0.5)
            # 触发查询
            page.evaluate("doQuery()")
            # 等待 KPI 数据加载
            try:
                page.wait_for_function(
                    "document.getElementById('kpiSchoolCount') && document.getElementById('kpiSchoolCount').textContent !== '--'",
                    timeout=20000,
                )
                time.sleep(3)  # 等图表渲染完成
            except Exception:
                time.sleep(5)
            kpi_path = str(tmp_dir / "home_kpi.png")
            try:
                el = page.query_selector("#kpiGrid")
                if el:
                    el.scroll_into_view_if_needed()
                    time.sleep(0.5)
                    el.screenshot(path=kpi_path)
                    screenshots["kpi"] = kpi_path
                    logger.info("[PPT] KPI 截图成功")
            except Exception as e:
                logger.warning("[PPT] KPI 截图失败: %s", e)
            
            # 滚动页面触发图表懒加载
            page.evaluate("window.scrollTo(0, document.body.scrollHeight)")
            time.sleep(2)  # 等待所有图表渲染
            page.evaluate("window.scrollTo(0, 0)")
            time.sleep(1)

            # ── 2. 首页柱状图（使用率排名） ──
            bar_path = str(tmp_dir / "home_bar.png")
            try:
                # 检查canvas是否存在
                canvas_exists = page.evaluate("!!document.getElementById('chartBar')")
                logger.info(f"[PPT] chartBar exists: {canvas_exists}")
                if canvas_exists:
                    el = page.query_selector("#chartGrid .chart-card:has(#chartBar)")
                    if not el:
                        el = page.query_selector("#chartGrid .chart-card:nth-child(2)")
                    if el:
                        el.scroll_into_view_if_needed()
                        rendered = _wait_for_canvas_rendered(page, "chartBar", timeout=15000)
                        logger.info(f"[PPT] chartBar rendered: {rendered}")
                        time.sleep(1)
                        el.screenshot(path=bar_path)
                        screenshots["bar"] = bar_path
                        logger.info("[PPT] 柱状图截图成功")
                    else:
                        logger.warning("[PPT] 未找到柱状图容器")
                else:
                    logger.warning("[PPT] chartBar canvas 不存在")
            except Exception as e:
                logger.warning("[PPT] 柱状图截图失败: %s", e)

            # ── 3. 环形图 + 雷达图（合并）──
            donut_radar_path = str(tmp_dir / "home_donut_radar.png")
            try:
                el = page.query_selector("#chartGrid")
                if el:
                    el.screenshot(path=donut_radar_path)
                    # 覆盖 bar，用整张图表区
                    screenshots["donut_radar"] = donut_radar_path
            except Exception as e:
                logger.warning("[PPT] 环形图截图失败: %s", e)

            # ── 4. 活跃比例对比图 ──
            ar_path = str(tmp_dir / "home_active_ratio.png")
            try:
                canvas_exists = page.evaluate("!!document.getElementById('chartActiveRatio')")
                logger.info(f"[PPT] chartActiveRatio exists: {canvas_exists}")
                if canvas_exists:
                    ar_el = page.query_selector("#chartGrid .chart-card:has(#chartActiveRatio)")
                    if not ar_el:
                        ar_el = page.query_selector("#chartGrid .chart-card:first-child")
                    if ar_el:
                        ar_el.scroll_into_view_if_needed()
                        rendered = _wait_for_canvas_rendered(page, "chartActiveRatio", timeout=15000)
                        logger.info(f"[PPT] chartActiveRatio rendered: {rendered}")
                        time.sleep(1)
                        ar_el.screenshot(path=ar_path)
                        screenshots["active_ratio"] = ar_path
                        logger.info("[PPT] 活跃比例图截图成功")
                    else:
                        logger.warning("[PPT] 未找到活跃比例图容器")
                else:
                    logger.warning("[PPT] chartActiveRatio canvas 不存在")
            except Exception as e:
                logger.warning("[PPT] 活跃比例图截图失败: %s", e)

            # ── 5. 作业次数统计图 ──
            hw_path = str(tmp_dir / "home_homework.png")
            try:
                # 最后一个 chart-card（各校作业次数统计）
                hw_el = page.query_selector("#chartGrid .chart-card:last-child")
                if hw_el:
                    hw_el.screenshot(path=hw_path)
                    screenshots["homework"] = hw_path
            except Exception as e:
                logger.warning("[PPT] 作业次数图截图失败: %s", e)

            # ── 5. 单校详情页 ──
            if school_id:
                school_url = f"{base_url}/school/{school_id}?start_date={start_date}&end_date={end_date}"
                logger.info("[PPT] 采集单校: %s", school_url)
                page.goto(school_url, wait_until="load", timeout=30000)
                time.sleep(1)
                try:
                    page.wait_for_function(
                        "document.getElementById('kpiUsage') && document.getElementById('kpiUsage').textContent !== '--'",
                        timeout=15000,
                    )
                    time.sleep(1.5)
                except Exception:
                    time.sleep(3)

                # 学校 KPI 卡片
                skpi_path = str(tmp_dir / "school_kpi.png")
                try:
                    el = page.query_selector(".school-kpi-grid")
                    if el:
                        el.screenshot(path=skpi_path)
                        screenshots["school_kpi"] = skpi_path
                except Exception as e:
                    logger.warning("[PPT] 单校KPI截图失败: %s", e)

                # KPI + 短板预警整块截图
                kpi_warn_path = str(tmp_dir / "school_kpi_warning.png")
                try:
                    # 确保短板预警可见
                    page.evaluate("""() => {
                        const warn = document.querySelector('.warning-section');
                        if (warn) warn.style.display = 'block';
                    }""")
                    time.sleep(0.5)
                    # 找到包含 KPI + warning 的容器并截取
                    container = page.query_selector(".school-kpi-grid")
                    if container:
                        # 获取 bounding box 并向上扩展以包含 warning section
                        box = container.bounding_box()
                        if box:
                            # 检查 warning section 是否存在并获取其位置
                            warn_box = page.evaluate("""() => {
                                const warn = document.querySelector('.warning-section');
                                if (warn) {
                                    const rect = warn.getBoundingClientRect();
                                    return { x: rect.x, y: rect.y, width: rect.width, height: rect.height, bottom: rect.bottom };
                                }
                                return null;
                            }""")
                            if warn_box:
                                # 扩展高度以包含 warning section
                                new_height = warn_box["bottom"] - box["y"]
                                page.screenshot(
                                    path=kpi_warn_path,
                                    clip={"x": box["x"], "y": box["y"], "width": box["width"], "height": new_height}
                                )
                                screenshots["school_kpi_warning"] = kpi_warn_path
                            else:
                                container.screenshot(path=kpi_warn_path)
                                screenshots["school_kpi_warning"] = kpi_warn_path
                except Exception as e:
                    logger.warning("[PPT] KPI+预警截图失败: %s", e)

                # ── 趋势分析 3 个 tab 分别截图 ──
                trend_section = page.query_selector(".trend-section")
                if trend_section:
                    trend_section.scroll_into_view_if_needed()
                    
                    # Tab 1: 使用率趋势（默认显示）
                    t1_path = str(tmp_dir / "school_trend_usage.png")
                    try:
                        _wait_for_canvas_rendered(page, "trendUsageChart")
                        time.sleep(0.5)
                        trend_section.screenshot(path=t1_path)
                        screenshots["school_trend_usage"] = t1_path
                        logger.info("[PPT] 使用率趋势截图成功")
                    except Exception as e:
                        logger.warning("[PPT] 使用率趋势截图失败: %s", e)

                    # Tab 2: 活跃度趋势
                    try:
                        page.click(".trend-tab:nth-child(2)")
                        _wait_for_canvas_rendered(page, "trendActivityChart")
                        time.sleep(0.5)
                        t2_path = str(tmp_dir / "school_trend_activity.png")
                        trend_section.screenshot(path=t2_path)
                        screenshots["school_trend_activity"] = t2_path
                        logger.info("[PPT] 活跃度趋势截图成功")
                    except Exception as e:
                        logger.warning("[PPT] 活跃度趋势截图失败: %s", e)

                    # Tab 3: 作业趋势
                    try:
                        page.click(".trend-tab:nth-child(3)")
                        _wait_for_canvas_rendered(page, "trendHomeworkChart")
                        time.sleep(0.5)
                        t3_path = str(tmp_dir / "school_trend_homework.png")
                        trend_section.screenshot(path=t3_path)
                        screenshots["school_trend_homework"] = t3_path
                        logger.info("[PPT] 作业趋势截图成功")
                    except Exception as e:
                        logger.warning("[PPT] 作业趋势截图失败: %s", e)

            browser.close()
    except Exception as e:
        logger.error("[PPT] 截图过程出错: %s", e, exc_info=True)

    return screenshots, tmp_dir


def _img_to_base64(img_path, max_width=400):
    """将图片转为 base64 缩略图"""
    try:
        from PIL import Image
        img = Image.open(img_path)
        w, h = img.size
        if w > max_width:
            ratio = max_width / w
            img = img.resize((max_width, int(h * ratio)), Image.LANCZOS)
        buf = BytesIO()
        img.save(buf, format="PNG")
        return base64.b64encode(buf.getvalue()).decode()
    except Exception:
        with open(img_path, "rb") as f:
            return base64.b64encode(f.read()).decode()


def _cleanup_session(session_id):
    """清理过期或已使用的会话"""
    info = _PPT_SESSIONS.pop(session_id, None)
    if info:
        try:
            shutil.rmtree(str(info["dir"]), ignore_errors=True)
        except Exception:
            pass


def _cleanup_expired():
    """清理所有过期会话"""
    now = time.time()
    expired = [sid for sid, info in _PPT_SESSIONS.items()
               if now - info["created_at"] > _PPT_SESSION_TTL]
    for sid in expired:
        _cleanup_session(sid)


@export_bp.route("/html-report", methods=["POST"])
def export_html_report():
    """导出独立 HTML 报告文件"""
    import json as json_mod
    data = request.get_json(silent=True) or {}
    start_date = data.get("start_date", "")
    end_date = data.get("end_date", "")
    school_id = data.get("school_id", "")
    report_title = data.get("report_title", "") or "教育平台数据看板"

    if not start_date or not end_date:
        return jsonify({"error": "缺少日期参数"}), 400

    # 调用内部 API 获取数据
    from flask import current_app
    client = current_app.test_client()
    with client.session_transaction() as sess:
        sess.update(session)

    # 1. 全校总览数据
    overview_resp = client.get(f"/api/charts/module-usage?start_date={start_date}&end_date={end_date}")
    overview_data = overview_resp.get_json() if overview_resp.status_code == 200 else {}

    # 2. 单校详情数据
    school_data = {}
    school_name = ""
    trend_data = {}
    if school_id:
        school_params = f"start_date={start_date}&end_date={end_date}&school_id={school_id}"
        school_resp = client.get(f"/api/charts/module-usage?{school_params}")
        school_data = school_resp.get_json() if school_resp.status_code == 200 else {}
        # 获取学校名称
        try:
            if school_data.get("rows"):
                school_name = school_data["rows"][0].get("display_name") or school_data["rows"][0].get("school", "")
            if not school_name:
                options_resp = client.get("/api/charts/options")
                if options_resp.status_code == 200:
                    opts = options_resp.get_json()
                    for schools in (opts.get("schools_by_type") or {}).values():
                        for s in schools:
                            if str(s.get("id")) == str(school_id):
                                school_name = s.get("display_name") or s.get("name", "")
                                break
        except Exception:
            pass
        # 趋势数据
        trend_resp = client.get(f"/api/charts/trend?{school_params}")
        trend_data = trend_resp.get_json() if trend_resp.status_code == 200 else {}

    html_content = _build_html_report(
        report_title, start_date, end_date, school_name,
        overview_data, school_data, trend_data
    )

    buf = io.BytesIO(html_content.encode("utf-8"))
    filename = f"数据报告_{start_date}_{end_date}.html"
    return send_file(buf, as_attachment=True, download_name=filename, mimetype="text/html")


def _build_html_report(title, start_date, end_date, school_name,
                       overview_data, school_data, trend_data):
    """生成自包含 HTML 报告"""
    import json as json_mod

    overview_json = json_mod.dumps(overview_data, ensure_ascii=False)
    school_json = json_mod.dumps(school_data, ensure_ascii=False)
    trend_json = json_mod.dumps(trend_data, ensure_ascii=False)

    # 提取 KPI
    rows = overview_data.get("rows", [])
    kpi_overview = {}
    if rows:
        avg_rate = sum((r.get("rate_values") or [0])[0] for r in rows) / len(rows)
        kpi_overview = {
            "school_count": len(rows),
            "avg_usage": round(avg_rate, 2),
            "daily_uv": sum(r.get("d21_uv", 0) or 0 for r in rows),
            "weekly_active": sum(r.get("d21_weekly_active", 0) or 0 for r in rows),
            "monthly_active": sum(r.get("d21_monthly_active", 0) or 0 for r in rows),
            "homework_total": sum((r.get("rate_values") or [0]*13)[8] or 0 for r in rows),
        }

    school_kpi = {}
    if school_data.get("rows"):
        s = school_data["rows"][0]
        rv = s.get("rate_values") or [0]*13
        school_kpi = {
            "usage_rate": rv[0] if len(rv) > 0 else 0,
            "daily_uv": s.get("d21_uv", 0) or 0,
            "weekly_active": s.get("d21_weekly_active", 0) or 0,
            "monthly_active": s.get("d21_monthly_active", 0) or 0,
            "homework_count": rv[8] if len(rv) > 8 else 0,
            "homework_per_capita": rv[9] if len(rv) > 9 else 0,
        }

    kpi_overview_json = json_mod.dumps(kpi_overview, ensure_ascii=False)
    school_kpi_json = json_mod.dumps(school_kpi, ensure_ascii=False)

    school_section = ""
    if school_name:
        school_section = f"""
<div class="section">
  <div class="section-title">{school_name} - 详细数据</div>
  <div class="kpi-grid" id="schoolKPI"></div>
  <div class="chart-card">
    <div class="title">周期趋势分析</div>
    <div class="tabs" id="trendTabs">
      <div class="tab active" data-tab="usage">使用率趋势</div>
      <div class="tab" data-tab="activity">活跃度趋势</div>
      <div class="tab" data-tab="homework">作业趋势</div>
    </div>
    <div class="tab-content active" id="tab-usage"><div class="chart-container"><canvas id="trendUsageChart"></canvas></div></div>
    <div class="tab-content" id="tab-activity"><div class="chart-container"><canvas id="trendActivityChart"></canvas></div></div>
    <div class="tab-content" id="tab-homework"><div class="chart-container"><canvas id="trendHomeworkChart"></canvas></div></div>
  </div>
</div>"""

    school_header = f'<div class="school-name">{school_name}</div>' if school_name else ''

    html = f'''<!DOCTYPE html>
<html lang="zh-CN">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>{title} - {start_date} ~ {end_date}</title>
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.4/dist/chart.umd.min.js"></script>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{font-family:-apple-system,BlinkMacSystemFont,"Segoe UI","Noto Sans SC",sans-serif;background:#f0f2f5;color:#1e293b;line-height:1.6}}
.container{{max-width:1200px;margin:0 auto;padding:20px}}
.header{{background:linear-gradient(135deg,#1e40af,#3b82f6);color:#fff;padding:40px;border-radius:16px;margin-bottom:24px;text-align:center}}
.header h1{{font-size:28px;margin-bottom:8px}}
.header .date{{font-size:15px;opacity:.85}}
.header .school-name{{font-size:18px;margin-top:8px;opacity:.9}}
.section{{margin-bottom:24px}}
.section-title{{font-size:20px;font-weight:600;color:#1e40af;margin-bottom:16px;padding-left:12px;border-left:4px solid #3b82f6}}
.kpi-grid{{display:grid;grid-template-columns:repeat(auto-fit,minmax(180px,1fr));gap:16px;margin-bottom:24px}}
.kpi-card{{background:#fff;border-radius:12px;padding:20px;text-align:center;box-shadow:0 1px 3px rgba(0,0,0,.1)}}
.kpi-card .value{{font-size:32px;font-weight:700;color:#1e40af}}
.kpi-card .label{{font-size:13px;color:#64748b;margin-top:4px}}
.chart-card{{background:#fff;border-radius:12px;padding:24px;margin-bottom:20px;box-shadow:0 1px 3px rgba(0,0,0,.1)}}
.chart-card .title{{font-size:16px;font-weight:600;margin-bottom:4px}}
.chart-card .subtitle{{font-size:13px;color:#94a3b8;margin-bottom:16px}}
.chart-container{{position:relative;height:360px}}
.chart-summary{{padding:12px 0;font-size:13px;color:#64748b;line-height:1.8;border-top:1px solid #e5e7eb;margin-top:12px}}
.tabs{{display:flex;gap:0;margin-bottom:20px}}
.tab{{padding:10px 24px;cursor:pointer;border:1px solid #d1d5db;background:#f9fafb;font-size:14px;color:#475569;transition:all .2s}}
.tab:first-child{{border-radius:8px 0 0 8px}}
.tab:last-child{{border-radius:0 8px 8px 0}}
.tab.active{{background:#3b82f6;color:#fff;border-color:#3b82f6}}
.tab-content{{display:none}}
.tab-content.active{{display:block}}
.footer{{text-align:center;padding:24px;color:#94a3b8;font-size:12px}}
@media print{{.container{{max-width:100%}}.chart-container{{height:280px}}}}
</style>
</head>
<body>
<div class="container">
<div class="header">
  <h1>{title}</h1>
  <div class="date">{start_date} ~ {end_date}</div>
  {school_header}
</div>
<div class="section">
  <div class="section-title">全校总览</div>
  <div class="kpi-grid" id="overviewKPI"></div>
  <div class="chart-card">
    <div class="title">学校活跃比例对比</div>
    <div class="subtitle">日活比例（折线）与周活比例（柱状）对比分析</div>
    <div class="chart-container"><canvas id="chartActiveRatio"></canvas></div>
    <div class="chart-summary" id="activeRatioSummary"></div>
  </div>
  <div class="chart-card">
    <div class="title">各校使用率排名</div>
    <div class="subtitle">筛选时间范围内的使用率排名</div>
    <div class="chart-container"><canvas id="chartBar"></canvas></div>
    <div class="chart-summary" id="barChartSummary"></div>
  </div>
  <div class="chart-card">
    <div class="title">各校作业次数统计</div>
    <div class="subtitle">筛选时间范围内的作业场次排名</div>
    <div class="chart-container"><canvas id="chartHomework"></canvas></div>
    <div class="chart-summary" id="homeworkSummary"></div>
  </div>
</div>
{school_section}
<div class="footer">数据来源：立达AI教学平台 · 报告生成时间 {datetime.now().strftime("%Y-%m-%d %H:%M")}</div>
</div>
<script>
var overviewData = {overview_json};
var schoolData = {school_json};
var trendData = {trend_json};
var kpiOverview = {kpi_overview_json};
var schoolKPI = {school_kpi_json};

function renderOverviewKPI() {{
  var el = document.getElementById("overviewKPI");
  if (!el) return;
  var items = [
    {{v: kpiOverview.school_count || 0, l: "学校总数"}},
    {{v: (kpiOverview.avg_usage || 0).toFixed(1) + "%", l: "平均使用率"}},
    {{v: kpiOverview.daily_uv || 0, l: "总日活人数"}},
    {{v: kpiOverview.weekly_active || 0, l: "总周活人数"}},
    {{v: kpiOverview.monthly_active || 0, l: "总月活人数"}},
    {{v: kpiOverview.homework_total || 0, l: "总作业次数"}}
  ];
  el.innerHTML = items.map(function(i){{return '<div class="kpi-card"><div class="value">'+i.v+'</div><div class="label">'+i.l+'</div></div>';}}).join("");
}}

function renderSchoolKPI() {{
  var el = document.getElementById("schoolKPI");
  if (!el) return;
  var items = [
    {{v: (schoolKPI.usage_rate || 0).toFixed(1) + "%", l: "平台使用率"}},
    {{v: schoolKPI.daily_uv || 0, l: "日活跃人数"}},
    {{v: schoolKPI.weekly_active || 0, l: "周活跃人数"}},
    {{v: schoolKPI.monthly_active || 0, l: "月活跃人数"}},
    {{v: (schoolKPI.homework_count||0) + "次 人均" + (schoolKPI.homework_per_capita||0).toFixed(1), l: "累计作业次数/人均"}}
  ];
  el.innerHTML = items.map(function(i){{return '<div class="kpi-card"><div class="value">'+i.v+'</div><div class="label">'+i.l+'</div></div>';}}).join("");
}}

function renderOverviewCharts() {{
  var rows = (overviewData.rows || []).slice().sort(function(a,b){{return ((b.rate_values||[])[0]||0)-((a.rate_values||[])[0]||0);}});
  if (!rows.length) return;
  var labels = rows.map(function(r){{return r.display_name;}});
  var dailyData = rows.map(function(r){{return (r.rate_values||[])[10]||0;}});
  var weeklyData = rows.map(function(r){{return (r.rate_values||[])[11]||0;}});
  new Chart(document.getElementById("chartActiveRatio"), {{
    type: "bar",
    data: {{
      labels: labels,
      datasets: [
        {{type:"line",label:"日活比例(%)",data:dailyData,borderColor:"#3b82f6",backgroundColor:"rgba(59,130,246,.1)",yAxisID:"y",tension:.3,pointRadius:3,borderWidth:2}},
        {{type:"bar",label:"周活比例(%)",data:weeklyData,backgroundColor:"rgba(34,197,94,.6)",yAxisID:"y1",borderRadius:3}}
      ]
    }},
    options: {{responsive:true,maintainAspectRatio:false,plugins:{{legend:{{position:"top"}}}},scales:{{
      y:{{type:"linear",position:"left",title:{{display:true,text:"日活比例(%)"}},beginAtZero:true}},
      y1:{{type:"linear",position:"right",title:{{display:true,text:"周活比例(%)"}},beginAtZero:true,grid:{{drawOnChartArea:false}}}},
      x:{{ticks:{{maxRotation:60,font:{{size:10}}}}}}
    }}}}
  }});
  var usageData = rows.map(function(r){{return (r.rate_values||[])[0]||0;}});
  var barColors = usageData.map(function(v){{return v>=30?"rgba(34,197,94,.7)":v>=15?"rgba(251,191,36,.7)":"rgba(239,68,68,.6)";}});
  new Chart(document.getElementById("chartBar"), {{
    type:"bar",
    data:{{labels:labels,datasets:[{{label:"使用率(%)",data:usageData,backgroundColor:barColors,borderRadius:3}}]}},
    options:{{responsive:true,maintainAspectRatio:false,plugins:{{legend:{{display:false}}}},scales:{{y:{{beginAtZero:true,title:{{display:true,text:"使用率(%)"}},ticks:{{callback:function(v){{return v+"%"}}}}}},x:{{ticks:{{maxRotation:60,font:{{size:10}}}}}}}}}}
  }});
  var sorted = rows.slice().sort(function(a,b){{return ((b.rate_values||[])[0]||0)-((a.rate_values||[])[0]||0);}});
  var top5 = sorted.slice(0,5).map(function(r){{return r.display_name+"("+(((r.rate_values||[])[0]||0).toFixed(1))+"%)";}}).join("、");
  var bot3 = sorted.slice(-3).reverse().map(function(r){{return r.display_name+"("+(((r.rate_values||[])[0]||0).toFixed(1))+"%)";}}).join("、");
  document.getElementById("barChartSummary").innerHTML = '<span style="color:#16a34a">使用率较高：'+top5+'</span><br><span style="color:#dc2626">使用率较低：'+bot3+'</span>';
  var hwRows = rows.map(function(r){{var rv=r.rate_values||[];return {{name:r.display_name,count:rv[8]||0,perCapita:rv[9]||0}};}}).filter(function(s){{return s.count>0;}}).sort(function(a,b){{return b.count-a.count;}});
  if (hwRows.length) {{
    new Chart(document.getElementById("chartHomework"), {{
      type:"bar",
      data:{{labels:hwRows.map(function(s){{return s.name;}}),datasets:[{{label:"作业次数",data:hwRows.map(function(s){{return s.count;}}),backgroundColor:"rgba(251,191,36,.6)",borderRadius:4}}]}},
      options:{{responsive:true,maintainAspectRatio:false,plugins:{{legend:{{display:false}},tooltip:{{callbacks:{{label:function(ctx){{var p=hwRows[ctx.dataIndex];return "作业次数: "+p.count+"次 (人均: "+p.perCapita.toFixed(1)+")";}}}}}}}},scales:{{y:{{beginAtZero:true,title:{{display:true,text:"作业次数"}}}},x:{{ticks:{{maxRotation:60,font:{{size:10}}}}}}}}}}
    }});
    var top5hw = hwRows.slice(0,5).map(function(s){{return s.name+"("+s.count+"次/人均"+s.perCapita.toFixed(1)+")";}}).join("、");
    document.getElementById("homeworkSummary").innerHTML = '使用较好的学校：<b style="color:#d97706">'+top5hw+'</b>';
  }}
}}

var _trendCharts = {{}};
function renderTrendCharts() {{
  if (!trendData || !trendData.labels) return;
  var labels = trendData.labels;
  var ds = trendData.datasets || [];
  var colorMap = {{"#22c55e":"rgba(34,197,94,1)","#3b82f6":"rgba(59,130,246,1)","#f97316":"rgba(249,115,22,1)","#a855f7":"rgba(168,85,247,1)"}};
  var colors = ["#22c55e","#3b82f6","#f97316","#a855f7"];
  var usageDS = ds.filter(function(d){{return d.label && d.label.indexOf("使用率")>=0;}});
  if (usageDS.length) {{
    _trendCharts.usage = new Chart(document.getElementById("trendUsageChart"), {{
      type:"line",
      data:{{labels:labels,datasets:usageDS.map(function(d,i){{return {{label:d.label,data:d.data,borderColor:colors[i%4],backgroundColor:"transparent",tension:.3,pointRadius:2,borderWidth:2}};}})}},
      options:{{responsive:true,maintainAspectRatio:false,plugins:{{legend:{{position:"bottom"}}}},scales:{{y:{{beginAtZero:true,ticks:{{callback:function(v){{return v+"%"}}}}}}}}}}
    }});
  }}
  var actDS = ds.filter(function(d){{return d.label && (d.label.indexOf("日活")>=0||d.label.indexOf("周活")>=0||d.label.indexOf("活跃")>=0) && d.label.indexOf("使用率")<0;}});
  if (actDS.length) {{
    _trendCharts.activity = new Chart(document.getElementById("trendActivityChart"), {{
      type:"line",
      data:{{labels:labels,datasets:actDS.map(function(d,i){{return {{label:d.label,data:d.data,borderColor:colors[i%4],backgroundColor:"transparent",tension:.3,pointRadius:2,borderWidth:2}};}})}},
      options:{{responsive:true,maintainAspectRatio:false,plugins:{{legend:{{position:"bottom"}}}},scales:{{y:{{beginAtZero:true}}}}}}
    }});
  }}
  var hwDS = ds.filter(function(d){{return d.label && d.label.indexOf("作业")>=0;}});
  if (hwDS.length) {{
    _trendCharts.homework = new Chart(document.getElementById("trendHomeworkChart"), {{
      type:"bar",
      data:{{labels:labels,datasets:hwDS.map(function(d){{return {{label:d.label,data:d.data,backgroundColor:"rgba(251,191,36,.6)",borderRadius:3}};}})}},
      options:{{responsive:true,maintainAspectRatio:false,plugins:{{legend:{{position:"bottom"}}}},scales:{{y:{{beginAtZero:true}}}}}}
    }});
  }}
}}

document.querySelectorAll(".tab").forEach(function(tab) {{
  tab.addEventListener("click", function() {{
    document.querySelectorAll(".tab").forEach(function(t){{t.classList.remove("active");}});
    document.querySelectorAll(".tab-content").forEach(function(c){{c.classList.remove("active");}});
    tab.classList.add("active");
    document.getElementById("tab-" + tab.getAttribute("data-tab")).classList.add("active");
    var key = tab.getAttribute("data-tab");
    if (_trendCharts[key]) setTimeout(function(){{_trendCharts[key].resize();}}, 100);
  }});
}});

renderOverviewKPI();
renderOverviewCharts();
renderSchoolKPI();
renderTrendCharts();
</script>
</body>
</html>'''
    return html


@export_bp.route("/ppt-screenshots", methods=["POST"])
def ppt_screenshots():
    """Step 1: 截取 PPT 所需截图，返回预览数据"""
    data = request.get_json(silent=True) or {}
    start_date = data.get("start_date", "") or request.args.get("start_date", "")
    end_date = data.get("end_date", "") or request.args.get("end_date", "")
    school_id = data.get("school_id", "") or request.args.get("school_id", "")
    ppt_title = data.get("ppt_title", "") or request.args.get("ppt_title", "")

    if not start_date or not end_date:
        return jsonify({"error": "时间范围为必填项"}), 400
    if not school_id:
        return jsonify({"error": "请指定学校"}), 400

    # 获取学校名称
    school_name = ""
    try:
        from models.school import School as LocalSchool
        s = LocalSchool.get_by_id(school_id)
        if s:
            school_name = s.name or s.school_name or ""
    except Exception:
        pass

    # 清理过期会话
    _cleanup_expired()

    # 创建会话
    session_id = uuid.uuid4().hex[:12]
    dest_dir = Path(tempfile.gettempdir()) / f"ppt_export_{session_id}"
    dest_dir.mkdir(parents=True, exist_ok=True)

    port = os.environ.get("PORT", "5001")
    base_url = f"http://127.0.0.1:{port}"

    logger.info("[PPT-Preview] 开始截屏: school=%s, %s ~ %s, session=%s",
                school_id, start_date, end_date, session_id)

    screenshots, tmp_dir = _capture_screenshots(
        base_url, start_date, end_date, school_id, school_name,
        dest_dir=str(dest_dir),
    )

    if not screenshots:
        shutil.rmtree(str(dest_dir), ignore_errors=True)
        return jsonify({"error": "截图失败，请检查服务是否正常运行"}), 500

    # 保存会话信息
    _PPT_SESSIONS[session_id] = {
        "dir": dest_dir,
        "start_date": start_date,
        "end_date": end_date,
        "school_name": school_name,
        "ppt_title": ppt_title,
        "created_at": time.time(),
    }

    # 构建预览数据
    slide_order = [
        ("cover", "封面", "天立国际 · 数据看板"),
    ]
    slide_keys_def = [
        ("kpi", "整体使用情况", "核心KPI指标"),
        ("active_ratio", "整体使用情况", "学校活跃比例对比"),
        ("bar", "整体使用情况", "各校使用率排名"),
        ("homework", "整体使用情况", "各校作业次数统计"),
    ]
    for key, section, label in slide_keys_def:
        if screenshots.get(key):
            slide_order.append((key, f"{section} — {label}", _PPT_DESC.get(key, "")))

    if school_name:
        school_slides = [
            ("school_kpi_warning", f"本校 — {school_name}", "核心KPI与短板预警"),
            ("school_trend_usage", f"本校 — {school_name}", "使用率趋势"),
            ("school_trend_activity", f"本校 — {school_name}", "活跃度趋势"),
            ("school_trend_homework", f"本校 — {school_name}", "作业趋势"),
        ]
        for key, section, label in school_slides:
            if screenshots.get(key):
                slide_order.append((key, f"{section} — {label}", _PPT_DESC.get(key, "")))

    slides = []
    for key, title, desc in slide_order:
        slide_data = {"key": key, "title": title, "desc": desc}
        if key == "cover":
            slide_data["isCover"] = True
        else:
            img_path = screenshots.get(key)
            if img_path and os.path.exists(img_path):
                slide_data["preview"] = _img_to_base64(img_path)
        slides.append(slide_data)

    return jsonify({
        "session_id": session_id,
        "school_name": school_name,
        "start_date": start_date,
        "end_date": end_date,
        "slides": slides,
    })


@export_bp.route("/ppt-download/<session_id>")
def ppt_download(session_id):
    """Step 2: 根据 session_id 生成并下载 PPT 文件"""
    _cleanup_expired()

    info = _PPT_SESSIONS.get(session_id)
    if not info:
        return jsonify({"error": "会话已过期，请重新截屏"}), 404

    dest_dir = info["dir"]
    if not dest_dir.exists():
        _cleanup_session(session_id)
        return jsonify({"error": "临时文件已丢失，请重新截屏"}), 404

    # 重新读取截图
    screenshots = {}
    for key in ["kpi", "bar", "active_ratio", "homework", "school_kpi", "school_kpi_warning",
                  "school_trend_usage", "school_trend_activity", "school_trend_homework",
                  "school_warning"]:
        candidates = [
            dest_dir / f"home_{key}.png",
            dest_dir / f"{key}.png",
        ]
        for c in candidates:
            if c.exists():
                screenshots[key] = str(c)
                break

    start_date = info["start_date"]
    end_date = info["end_date"]
    school_name = info.get("school_name", "")
    ppt_title = info.get("ppt_title", "")

    logger.info("[PPT-Download] 生成 PPT: session=%s, school=%s, title=%s", session_id, school_name, ppt_title)

    try:
        pptx_buf = _make_pptx(screenshots, start_date, end_date, school_name, ppt_title=ppt_title)
    except Exception as e:
        logger.error("[PPT] PPT 生成失败: %s", e, exc_info=True)
        return jsonify({"error": f"PPT 生成失败: {str(e)}"}), 500

    # 清理
    _cleanup_session(session_id)

    dl_name = f"数据看板_{start_date}_{end_date}_{school_name or 'export'}.pptx"
    return send_file(
        pptx_buf,
        as_attachment=True,
        download_name=dl_name,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
